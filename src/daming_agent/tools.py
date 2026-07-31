import base64
import hashlib
import json
import os
import subprocess
import threading
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Optional, Any, Union

import httpx
from bs4 import BeautifulSoup
try:
    from ddgs import DDGS
except ImportError:
    from duckduckgo_search import DDGS
from playwright.sync_api import sync_playwright, Playwright, BrowserContext, Page
from runtime_store import RuntimeStore


@dataclass
class BrowserSession:
    """浏览器状态必须绑定到一个会话，禁止渠道之间复用页面或登录态。"""
    context: BrowserContext
    page: Page
    headless: bool
    read_only_view: bool


class TaskManager:
    """本地后台异步任务管理器。"""

    def __init__(self, workspace: Path, store: RuntimeStore | None = None, session_id: str = "default") -> None:
        self.workspace = workspace
        self.session_id = session_id
        self.tasks: dict[str, dict[str, Any]] = {}
        self.store = store or RuntimeStore(workspace / ".runtime")
        self._recover_tasks()

    def _recover_tasks(self) -> None:
        """重启后绝不重跑任务；仅记录仍存活 PID 为待人工恢复监控。"""
        for record in self.store.tasks(self.session_id, limit=500):
            if record.get("workspace") != str(self.workspace) or record.get("status") not in {"running", "running_recovered"}:
                continue
            pid = record.get("pid")
            alive = isinstance(pid, int)
            if alive:
                try: os.kill(pid, 0)
                except OSError: alive = False
            self.store.task(record["id"], status="running_recovered" if alive else "interrupted", recovery_note="进程仍存活，需重新附加监控" if alive else "Agent 重启后无法安全接管，未自动重跑")

    def run_command_async(self, command: str, relative_cwd: str = ".", objective: str = "") -> str:
        task_id = f"task_{uuid.uuid4().hex[:12]}"
        cwd_path = (self.workspace / (relative_cwd or ".")).resolve()

        log_path = self.workspace / ".logs" / f"{task_id}.log"
        log_path.parent.mkdir(parents=True, exist_ok=True)
        log_file = log_path.open("w", encoding="utf-8")

        process = subprocess.Popen(
            command,
            shell=True,
            cwd=str(cwd_path),
            stdout=log_file,
            stderr=subprocess.STDOUT,
            text=True,
        )
        self.tasks[task_id] = {
            "id": task_id,
            "command": command,
            "pid": process.pid,
            "process": process,
            "log_file": log_file,
            "log_path": log_path,
            "start_time": time.strftime("%Y-%m-%d %H:%M:%S"),
        }
        self.store.task(task_id, session_id=self.session_id, command=command, objective=objective, pid=process.pid, status="running", progress="已启动", log_path=str(log_path), workspace=str(self.workspace), retry_count=0, artifacts=[str(log_path)])
        return f"已在后台成功启动任务 [{task_id}] (PID: {process.pid})。日志输出在: {log_path.relative_to(self.workspace)}"

    def list_tasks(self) -> str:
        rows = self.store.tasks(self.session_id, limit=100)
        if not rows:
            return "当前没有在后台运行或已记录的任务。"
        lines = []
        for record in rows:
            tid, task = record["id"], self.tasks.get(record["id"])
            if task:
                exit_code = task["process"].poll()
                if exit_code is not None:
                    log_file = task.get("log_file")
                    if log_file and not log_file.closed: log_file.close()
                    self.store.task(tid, status="completed" if exit_code == 0 else "failed", exit_code=exit_code, progress="已结束")
                    record = self.store.get_task(tid, self.session_id) or record
            lines.append(f"[{tid}] 命令: '{record.get('command', '')}' | PID: {record.get('pid', '-')} | 状态: {record.get('status')} | 重试: {record.get('retry_count', 0)}")
        return "\n".join(lines)

    def get_task_log(self, task_id: str, tail_lines: int = 50) -> str:
        record = self.store.get_task(task_id, self.session_id)
        if not record:
            return f"任务不存在: {task_id}"
        log_path = Path(record.get("log_path", ""))
        if not log_path.exists():
            return "日志文件尚不存在。"
        content = log_path.read_text(encoding="utf-8", errors="ignore")
        lines = content.splitlines()
        tail = lines[-tail_lines:] if len(lines) > tail_lines else lines
        return "\n".join(tail) or "日志内容为空。"

    def kill_task(self, task_id: str) -> str:
        record = self.store.get_task(task_id, self.session_id)
        if not record:
            return f"任务不存在: {task_id}"
        process = self.tasks.get(task_id, {}).get("process")
        if process is None:
            return f"任务 [{task_id}] 不能由当前 Agent 安全终止（重启后未接管）。"
        if process.poll() is not None:
            return f"任务 [{task_id}] 已经停止执行。"
        try:
            process.terminate()
            time.sleep(0.5)
            if process.poll() is None:
                process.kill()
            log_file = self.tasks.get(task_id, {}).get("log_file")
            if log_file and not log_file.closed: log_file.close()
            self.store.task(task_id, status="cancelled", progress="已取消")
            return f"已成功终止后台任务 [{task_id}]。"
        except Exception as e:
            return f"终止任务 [{task_id}] 失败: {e}"

    def retry_task(self, task_id: str) -> str:
        record = self.store.get_task(task_id, self.session_id)
        if not record: return f"任务不存在: {task_id}"
        if record.get("status") in {"running", "running_recovered"}: return f"任务 [{task_id}] 仍在运行，不能重试。"
        result = self.run_command_async(record.get("command", ""), ".", record.get("objective", ""))
        # 新任务是独立 id，保留重试链而不会覆盖原始审计。
        new_id = result.split("[")[1].split("]")[0] if "[" in result else ""
        if new_id: self.store.task(new_id, retry_count=int(record.get("retry_count", 0)) + 1, retry_of=task_id)
        return result


class LocalTools:
    """本地 Agent 工具集：包括操作 Workspace 内文件、网络搜索、网页提取、Playwright 浏览器与 Shell 命令。"""

    def __init__(self, workspace: Path, headless: bool = True, slow_mo_ms: int = 80, runtime_store: RuntimeStore | None = None, session_id: str = "default") -> None:
        self.workspace = workspace.resolve()
        self.workspace.mkdir(parents=True, exist_ok=True)
        self.task_manager = TaskManager(self.workspace, runtime_store, session_id)
        self._pw: Optional[Playwright] = None
        self._browser_sessions: dict[str, BrowserSession] = {}
        self._browser_lock = threading.RLock()
        self.browser_headless = bool(headless)
        self.browser_slow_mo_ms = max(0, int(slow_mo_ms))
        self._browser_profiles = self.workspace / ".browser_profiles"
        self._browser_profiles.mkdir(exist_ok=True)


    # --- 文件读取与列表 ---

    def list_files(self, relative_path: str) -> str:
        try:
            path = self._safe_path(relative_path or ".")
            if not path.exists():
                return "路径不存在。"
            if path.is_file():
                return path.name
            files = [str(item.relative_to(self.workspace)) for item in path.rglob("*") if item.is_file()]
            return "\n".join(files[:100]) or "文件夹为空。"
        except ValueError:
            return "拒绝访问：只能读取 workspace 文件夹内的内容。"

    def read_text_file(self, relative_path: str) -> str:
        try:
            path = self._safe_path(relative_path or "")
            if not path.is_file():
                return "文件不存在，或这不是一个文件。"
            if path.stat().st_size > 100_000:
                return "文件过大，当前版本只允许读取小于 100 KB 的文本文件。"
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return "这不是 UTF-8 文本文件，当前版本不读取它。"
        except ValueError:
            return "拒绝访问：只能读取 workspace 文件夹内的内容。"

    # --- 文件写入与编辑 (P1) ---

    def write_file(self, relative_path: str, content: str) -> str:
        try:
            path = self._safe_path(relative_path or "")
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_text(content, encoding="utf-8")
            return f"成功写入文件: {path.relative_to(self.workspace)}"
        except ValueError:
            return "拒绝访问：只能在 workspace 文件夹内写入文件。"
        except Exception as e:
            return f"写入文件失败: {e}"

    def append_file(self, relative_path: str, content: str) -> str:
        try:
            path = self._safe_path(relative_path or "")
            path.parent.mkdir(parents=True, exist_ok=True)
            with path.open("a", encoding="utf-8") as f:
                f.write(content)
            return f"成功追加内容到文件: {path.relative_to(self.workspace)}"
        except ValueError:
            return "拒绝访问：只能在 workspace 文件夹内编辑文件。"
        except Exception as e:
            return f"编辑文件失败: {e}"

    def read_file_lines(self, relative_path: str, start_line: int = 1, end_line: int = 100) -> str:
        """精准按指定行号范围读取文本文件。"""
        try:
            path = self._safe_path(relative_path or "")
            if not path.is_file():
                return "文件不存在，或这不是一个文件。"
            lines = path.read_text(encoding="utf-8").splitlines()
            start = max(1, start_line)
            end = min(len(lines), end_line)
            selected = lines[start - 1 : end]
            formatted = [f"{idx}: {line}" for idx, line in enumerate(selected, start=start)]
            return f"=== 文件 {path.relative_to(self.workspace)} (第 {start}-{end} 行，共 {len(lines)} 行) ===\n" + "\n".join(formatted)
        except Exception as e:
            return f"按行读取文件失败: {e}"

    def replace_file_content(self, relative_path: str, target_text: str, replacement_text: str, expected_hash: Optional[str] = None) -> str:
        """精准匹配并替换文本/代码块，支持可选的 expected_hash 哈希防漂移校验。"""
        try:
            path = self._safe_path(relative_path or "")
            if not path.is_file():
                return "文件不存在，或这不是一个文件。"
            content = path.read_text(encoding="utf-8")
            if target_text not in content:
                return "替换失败：未在文件中找到精准匹配的 Target 文本。请检查缩进与换行。"
            if expected_hash:
                import hashlib
                actual_hash = hashlib.md5(target_text.encode("utf-8")).hexdigest()[:8]
                if actual_hash.lower() != expected_hash.lower().strip():
                    return f"❌ 替换拒绝：Hash 签名校验失败（期望 {expected_hash}，实际 {actual_hash}）。目标内容已被修改或发生行号漂移，请重新读取文件后再编辑。"
            count = content.count(target_text)
            if count > 1:
                return f"替换失败：Target 文本在文件中出现了 {count} 次，不够唯一。请提供更长更独特的上下文。"
            new_content = content.replace(target_text, replacement_text)
            path.write_text(new_content, encoding="utf-8")
            return f"成功精准替换文件 {path.relative_to(self.workspace)} 中的文本内容。"
        except Exception as e:
            return f"替换文件内容失败: {e}"


    def move_file(self, source_path: str, destination_path: str) -> str:
        """重命名或移动 workspace 内部的文件或文件夹。"""
        try:
            src = self._safe_path(source_path)
            dst = self._safe_path(destination_path)
            if not src.exists():
                return f"源文件/文件夹不存在: {source_path}"
            dst.parent.mkdir(parents=True, exist_ok=True)
            src.rename(dst)
            return f"成功移动/重命名: {src.name} -> {dst.relative_to(self.workspace)}"
        except ValueError:
            return "拒绝操作：源或目标路径必须限定在 workspace 文件夹内。"
        except Exception as e:
            return f"移动/重命名失败: {e}"

    def delete_file(self, relative_path: str) -> str:
        """删除 workspace 内部的指定文件。"""
        try:
            path = self._safe_path(relative_path)
            if not path.exists():
                return f"文件不存在: {relative_path}"
            if path.is_dir():
                import shutil
                shutil.rmtree(path)
                return f"成功删除文件夹及其全部内容: {path.relative_to(self.workspace)}"
            else:
                path.unlink()
                return f"成功删除文件: {path.relative_to(self.workspace)}"
        except ValueError:
            return "拒绝操作：只能删除 workspace 文件夹内部的内容。"
        except Exception as e:
            return f"删除文件失败: {e}"


    # --- 联网搜索与网页抓取 (P0) ---

    def web_search(self, query: str) -> str:
        if not query.strip():
            return "搜索关键词不能为空。"
        
        # 1. 引擎一：DuckDuckGo DDGS API
        try:
            with DDGS() as ddgs:
                results = list(ddgs.text(query, max_results=5))
            if results:
                formatted = []
                for i, r in enumerate(results, 1):
                    title = r.get("title", "无标题")
                    snippet = r.get("body", "")
                    link = r.get("href", "")
                    formatted.append(f"[{i}] {title}\n    网址: {link}\n    摘要: {snippet}")
                return "\n\n".join(formatted)
        except Exception as e:
            print(f"⚠️ [DDGS API 搜索降级引发]: {e}，切至备用 DuckDuckGo Lite HTML 抓取引擎...")

        # 2. 引擎二：DuckDuckGo Lite HTML 网页解构器
        try:
            headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"}
            resp = httpx.post("https://lite.duckduckgo.com/lite/", data={"q": query}, headers=headers, timeout=10.0)
            if resp.status_code == 200:
                soup = BeautifulSoup(resp.text, "html.parser")
                rows = soup.select(".result-snippet")
                links = soup.select(".result-link")
                results = []
                for i in range(min(5, len(links))):
                    title = links[i].get_text(strip=True)
                    link = links[i].get("href", "")
                    snippet = rows[i].get_text(strip=True) if i < len(rows) else ""
                    results.append(f"[{i+1}] {title}\n    网址: {link}\n    摘要: {snippet}")
                if results:
                    return "🌐 [使用备用 DDG Lite 引擎完成搜索]:\n\n" + "\n\n".join(results)
        except Exception as e:
            print(f"⚠️ [DDG Lite 搜索降级引发]: {e}，切至第三方 Bing 备用引擎...")

        # 3. 引擎三：Bing HTML 搜索兜底
        try:
            headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"}
            resp = httpx.get(f"https://www.bing.com/search?q={httpx.URL(query).raw_path.decode() if hasattr(httpx.URL(query), 'raw_path') else query}", headers=headers, timeout=10.0)
            if resp.status_code == 200:
                soup = BeautifulSoup(resp.text, "html.parser")
                items = soup.select("li.b_algo")
                results = []
                for i, item in enumerate(items[:5], 1):
                    a_tag = item.select_one("h2 a")
                    p_tag = item.select_one(".b_caption p")
                    if a_tag:
                        title = a_tag.get_text(strip=True)
                        link = a_tag.get("href", "")
                        snippet = p_tag.get_text(strip=True) if p_tag else ""
                        results.append(f"[{i}] {title}\n    网址: {link}\n    摘要: {snippet}")
                if results:
                    return "🌐 [使用 Bing 兜底引擎完成搜索]:\n\n" + "\n\n".join(results)
        except Exception as e:
            print(f"⚠️ [Bing 兜底引擎异常]: {e}")

        return "未找到相关搜索结果或所有搜索引擎不可用。"

    def get_token_stats(self) -> str:
        """获取当前 Agent 的全局 Token 消耗明细与统计报表。"""
        token_file = self.workspace.parent / "data" / "token_usage.json"
        if not token_file.exists():
            return "暂无 Token 消耗数据。"
        try:
            content = token_file.read_text(encoding="utf-8")
            data = json.loads(content)
            summary = [
                "📊 [Agent 全局 Token 消耗统计报表]",
                f"  - 累计请求次数: {data.get('total_requests', 0)} 次",
                f"  - 累计 Prompt (输入): {data.get('total_prompt_tokens', 0):,} tokens",
                f"  - 累计 Completion (输出): {data.get('total_completion_tokens', 0):,} tokens",
                f"  - 累计总用量: {data.get('total_tokens', 0):,} tokens\n",
                "模型消耗明细:"
            ]
            for model_name, stats in data.get("by_model", {}).items():
                summary.append(
                    f"  • [{model_name}]: {stats.get('total_tokens', 0):,} tokens ({stats.get('requests', 0)} 次请求)"
                )
            return "\n".join(summary)
        except Exception as e:
            return f"读取 Token 统计信息失败: {e}"

    def fetch_webpage(self, url: str) -> str:

        if not url.startswith(("http://", "https://")):
            return "URL 格式无效，必须以 http:// 或 https:// 开头。"
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            }
            response = httpx.get(url, headers=headers, follow_redirects=True, timeout=15.0)
            response.raise_for_status()

            soup = BeautifulSoup(response.text, "html.parser")
            for script_or_style in soup(["script", "style", "noscript", "header", "footer", "svg"]):
                script_or_style.decompose()

            text = soup.get_text(separator="\n", strip=True)
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            cleaned_text = "\n".join(lines)
            if len(cleaned_text) > 4000:
                cleaned_text = cleaned_text[:4000] + "\n...(内容过长已截断)"
            return cleaned_text or "网页内容为空。"
        except Exception as e:
            return f"抓取网页失败: {e}"

    # --- Playwright 后台浏览器自动化 (P1) ---

    def _get_browser_session(self, session_id: str) -> Optional[BrowserSession]:
        return self._browser_sessions.get(session_id)

    def _profile_dir(self, session_id: str) -> Path:
        """会话专属、可跨重启保存的登录态；路径不泄露会话身份。"""
        digest = hashlib.sha256(session_id.encode("utf-8")).hexdigest()
        return self._browser_profiles / digest

    @staticmethod
    def _viewer_init_script() -> str:
        """给可见页面加只读保护和 AI 光标图层，不改变站点业务 DOM。"""
        return """
(() => {
  window.__damingAgentAllowsInput = false;
  const blockUserInput = (event) => {
    if (!window.__damingAgentAllowsInput && event.isTrusted) {
      event.preventDefault(); event.stopImmediatePropagation();
    }
  };
  for (const name of ['pointerdown','pointerup','mousedown','mouseup','click','dblclick','keydown','keyup','keypress','paste','cut','drop','contextmenu']) {
    document.addEventListener(name, blockUserInput, true);
  }
  const mount = () => {
    if (document.getElementById('__daming_agent_cursor')) return;
    const cursor = document.createElement('div');
    cursor.id = '__daming_agent_cursor';
    cursor.setAttribute('aria-hidden', 'true');
    cursor.style.cssText = 'position:fixed;z-index:2147483647;left:-40px;top:-40px;width:22px;height:22px;border:3px solid #22d3ee;border-radius:50%;background:rgba(34,211,238,.18);pointer-events:none;transform:translate(-50%,-50%);transition:left .18s ease-out,top .18s ease-out,transform .18s ease-out;box-shadow:0 0 0 5px rgba(34,211,238,.16),0 2px 12px rgba(0,0,0,.45)';
    const label = document.createElement('span');
    label.style.cssText = 'position:absolute;left:15px;top:15px;white-space:nowrap;padding:3px 6px;border-radius:5px;background:#0f172a;color:#fff;font:12px system-ui;box-shadow:0 1px 5px rgba(0,0,0,.35)';
    label.textContent = 'AI 操作中'; cursor.appendChild(label); document.documentElement.appendChild(cursor);
    window.__damingAgentCursor = (x, y, labelText, click) => {
      cursor.style.left = x + 'px'; cursor.style.top = y + 'px'; label.textContent = labelText || 'AI 操作中';
      cursor.style.transform = click ? 'translate(-50%,-50%) scale(1.65)' : 'translate(-50%,-50%) scale(1)';
      if (click) setTimeout(() => { cursor.style.transform = 'translate(-50%,-50%) scale(1)'; }, 180);
    };
  };
  if (document.readyState === 'loading') document.addEventListener('DOMContentLoaded', mount, {once:true}); else mount();
})();
"""

    @staticmethod
    def _set_agent_input(page: Page, allowed: bool) -> None:
        page.evaluate("allowed => { window.__damingAgentAllowsInput = allowed; }", allowed)

    @staticmethod
    def _show_agent_cursor(page: Page, x: float, y: float, label: str, click: bool = False) -> None:
        page.evaluate("([x, y, label, click]) => window.__damingAgentCursor?.(x, y, label, click)", [x, y, label, click])

    def open_browser(self, url: str, session_id: str = "default", show_window: bool = False) -> str:
        if not url.startswith(("http://", "https://")):
            url = "https://" + url
        try:
            with self._browser_lock:
                browser_session = self._get_browser_session(session_id)
                requested_headless = not show_window and self.browser_headless
                if browser_session and show_window and browser_session.headless:
                    browser_session.context.close()
                    self._browser_sessions.pop(session_id, None)
                    browser_session = None
                if browser_session is None:
                    if self._pw is None:
                        self._pw = sync_playwright().start()
                    # 使用本机 Chrome 的独立 profile：兼容性优于 Playwright 自带
                    # Chromium，且登录态不会写入用户日常 Chrome profile。
                    launch_options = {
                        "headless": requested_headless,
                        "slow_mo": self.browser_slow_mo_ms,
                        "channel": "chrome",
                        "args": ["--test-type", "--disable-blink-features=AutomationControlled"],
                        "ignore_default_args": ["--enable-automation"],
                    }
                    try:
                        context = self._pw.chromium.launch_persistent_context(str(self._profile_dir(session_id)), **launch_options)
                    except Exception:
                        launch_options.pop("channel", None)
                        context = self._pw.chromium.launch_persistent_context(str(self._profile_dir(session_id)), **launch_options)
                    context.add_init_script(self._viewer_init_script())
                    page = context.pages[0] if context.pages else context.new_page()
                    browser_session = BrowserSession(context, page, requested_headless, read_only_view=not requested_headless)
                    self._browser_sessions[session_id] = browser_session

                self._set_agent_input(browser_session.page, True)
                try:
                    browser_session.page.goto(url, wait_until="domcontentloaded", timeout=20000)
                finally:
                    # 默认启用独占保护锁定，防止人类误触干扰 AI DOM 计算
                    self._set_agent_input(browser_session.page, False)
                title = browser_session.page.title()
                current_url = browser_session.page.url
                body_text = browser_session.page.inner_text("body")[:1000].replace("\n", " ")
                mode = "后台浏览器" if browser_session.headless else "可见观察窗口(人类操作已锁定)"
                return f"已在{mode}中打开网页:\n标题: {title}\n网址: {current_url}\n页面文本预览: {body_text}..."
        except Exception as e:
            return f"打开浏览器失败: {e}"

    def click(self, selector: str, session_id: str = "default") -> str:
        browser_session = self._get_browser_session(session_id)
        if browser_session is None:
            return "浏览器尚未打开，请先调用 open_browser。"
        try:
            with self._browser_lock:
                # 支持选择器或文本匹配，并把目标位置投射为可见 AI 光标轨迹。
                target = browser_session.page.locator(selector)
                if not selector.startswith((".", "#", "[", "text=")):
                    try:
                        target = browser_session.page.get_by_text(selector, exact=True)
                        target.bounding_box(timeout=5000)
                    except Exception:
                        target = browser_session.page.locator(selector)
                box = target.bounding_box(timeout=5000)
                if box:
                    self._show_agent_cursor(
                        browser_session.page,
                        box["x"] + box["width"] / 2,
                        box["y"] + box["height"] / 2,
                        "AI 点击",
                        click=True,
                    )
                self._set_agent_input(browser_session.page, True)
                try:
                    target.click(timeout=5000)
                finally:
                    self._set_agent_input(browser_session.page, False)
            return f"成功点击元素: {selector}"
        except Exception as e:
            return f"点击元素 {selector} 失败: {e}。提示：如属复杂/Canvas/无 DOM 元素，请使用 click_visual(description, x_ratio, y_ratio) 进行视觉物理坐标点击。"

    def click_visual(self, description: str = "", x_ratio: float = 0.5, y_ratio: float = 0.5, session_id: str = "default") -> str:
        """基于比例/物理坐标预测或描述执行视觉 Mouse 点击降级。"""
        browser_session = self._get_browser_session(session_id)
        if browser_session is None:
            return "浏览器尚未打开，请先调用 open_browser。"
        try:
            with self._browser_lock:
                page = browser_session.page
                viewport = page.viewport_size or {"width": 1280, "height": 720}
                x = viewport["width"] * max(0.0, min(1.0, float(x_ratio)))
                y = viewport["height"] * max(0.0, min(1.0, float(y_ratio)))
                self._show_agent_cursor(page, x, y, f"视觉点击: {description}", click=True)
                self._set_agent_input(page, True)
                try:
                    page.mouse.click(x, y)
                finally:
                    self._set_agent_input(page, False)
            return f"成功执行视觉 GUI 点击: [{description}] 坐标位置 ({x:.1f}, {y:.1f})"
        except Exception as e:
            return f"视觉 GUI 点击失败: {e}"

    def type_text(self, selector: str, text: str, session_id: str = "default") -> str:
        browser_session = self._get_browser_session(session_id)
        if browser_session is None:
            return "浏览器尚未打开，请先调用 open_browser。"
        try:
            with self._browser_lock:
                # 用逐字输入代替瞬间 fill，更接近日常浏览器操作。
                target = browser_session.page.locator(selector)
                box = target.bounding_box(timeout=5000)
                if box:
                    self._show_agent_cursor(
                        browser_session.page,
                        box["x"] + box["width"] / 2,
                        box["y"] + box["height"] / 2,
                        "AI 输入",
                    )
                self._set_agent_input(browser_session.page, True)
                try:
                    target.click(timeout=5000)
                    target.press("ControlOrMeta+A", timeout=5000)
                    target.press("Backspace", timeout=5000)
                    target.press_sequentially(text, delay=max(20, self.browser_slow_mo_ms))
                finally:
                    self._set_agent_input(browser_session.page, False)
            return f"成功在 {selector} 输入文本: {text}"
        except Exception as e:
            return f"输入文本失败: {e}"

    def press_key(self, key: str = "Enter", session_id: str = "default") -> str:
        """在浏览器当前页面模拟键盘按键 (如 Enter, Escape, Tab, Backspace, ArrowDown 等)。"""
        browser_session = self._get_browser_session(session_id)
        if browser_session is None:
            return "浏览器尚未打开，请先调用 open_browser。"
        try:
            with self._browser_lock:
                self._set_agent_input(browser_session.page, True)
                try:
                    browser_session.page.keyboard.press(key)
                finally:
                    self._set_agent_input(browser_session.page, not browser_session.headless)
            return f"已成功在浏览器按键: {key}"
        except Exception as e:
            return f"模拟按键 {key} 失败: {e}"

    def screenshot(self, relative_path: str = "screenshot.png", session_id: str = "default") -> str:
        browser_session = self._get_browser_session(session_id)
        if browser_session is None:
            return "浏览器尚未打开，请先调用 open_browser。"
        try:
            with self._browser_lock:
                path = self._safe_path(relative_path)
                path.parent.mkdir(parents=True, exist_ok=True)
                bytes_data = browser_session.page.screenshot(path=str(path))
            b64_str = base64.b64encode(bytes_data).decode("utf-8")
            result = {
                "type": "screenshot_result",
                "message": f"已保存页面截图到: {path.relative_to(self.workspace)}",
                "path": str(path.relative_to(self.workspace)),
                "base64": b64_str
            }
            return json.dumps(result, ensure_ascii=False)
        except Exception as e:
            return f"截图失败: {e}"


    def close_browser(self, session_id: str = "default") -> str:
        try:
            with self._browser_lock:
                browser_session = self._browser_sessions.pop(session_id, None)
                if browser_session is None:
                    return "当前会话没有打开浏览器。"
                browser_session.context.close()
                if not self._browser_sessions and self._pw:
                    self._pw.stop()
                    self._pw = None
                return "已成功关闭当前会话的浏览器。"
        except Exception as e:
            return f"关闭浏览器出错: {e}"

    def request_human_intervention(self, reason: str = "遇到需要人工干预的场景（如扫码登录、人机验证等）", session_id: str = "default") -> str:
        """当遇到登录框、二维码或人机校验时，临时解除物理保护锁定，允许人类手工操作。"""
        browser_session = self._get_browser_session(session_id)
        if browser_session is None:
            return "浏览器尚未打开，请先调用 open_browser。"
        try:
            with self._browser_lock:
                # 临时解禁人类物理操作
                self._set_agent_input(browser_session.page, True)
                self._show_agent_cursor(browser_session.page, 120, 30, f"⚠️ 等待人类操作: {reason[:30]}", click=False)
            return f"已临时开启人类操作权限（原因: {reason}）。请在窗口完成手工操作/登录后告诉我。"
        except Exception as e:
            return f"开启人类操作权限失败: {e}"

    def resume_agent_control(self, session_id: str = "default") -> str:
        """人类手工操作完成后，重新恢复 AI 独占保护锁定并恢复后续自动化流程。"""
        browser_session = self._get_browser_session(session_id)
        if browser_session is None:
            return "浏览器尚未打开。"
        try:
            with self._browser_lock:
                # 重新加锁保护
                self._set_agent_input(browser_session.page, False)
                self._show_agent_cursor(browser_session.page, -40, -40, "AI 已接管", click=False)
            return "AI 已重新接管控制权，保护锁定已恢复。"
        except Exception as e:
            return f"恢复 AI 保护锁定失败: {e}"

    # --- 安全 Shell 命令执行 (P2) ---

    def run_command(self, command: str, relative_cwd: str = ".") -> str:
        if not command.strip():
            return "命令不能为空。"
        try:
            cwd_path = self._safe_path(relative_cwd or ".")
            if not cwd_path.is_dir():
                return "工作目录不存在或不是文件夹。"

            # 挂载 SandboxGate 编译级静态 AST 语法安检门校验
            try:
                import ast
                parts = command.strip().split()
                if len(parts) >= 2 and parts[0] in {"python", "python3"} and parts[1].endswith(".py"):
                    script_file = cwd_path / parts[1]
                    if script_file.exists() and script_file.is_file():
                        code = script_file.read_text(encoding="utf-8", errors="ignore")
                        ast.parse(code)
            except SyntaxError as syntax_err:
                return f"⚠️ SandboxGate 静态语法安检拦截: 关联 Python 脚本语法报错 ({syntax_err})"
            except Exception:
                pass



            result = subprocess.run(
                command,
                shell=True,
                cwd=str(cwd_path),
                capture_output=True,
                text=True,
                timeout=30,
            )
            stdout = result.stdout.strip()
            stderr = result.stderr.strip()
            code = result.returncode

            output = [f"退出码: {code}"]
            if stdout:
                output.append(f"标准输出:\n{stdout}")
            if stderr:
                output.append(f"标准错误输出:\n{stderr}")
            return "\n".join(output)
        except subprocess.TimeoutExpired:
            return "命令执行超时（上限 30 秒）。"
        except ValueError:
            return "拒绝执行：工作目录只能限定在 workspace 内。"
        except Exception as e:
            return f"执行命令失败: {e}"

    # --- 办公文档读写与生成 (Office Documents) ---

    def read_office_file(self, relative_path: str) -> str:
        """读取 PDF, Word (.docx), PPT (.pptx), Excel (.xlsx/.csv) 文档纯文本。"""
        try:
            path = self._safe_path(relative_path)
            if not path.is_file():
                return "文件不存在，或这不是一个文件。"

            ext = path.suffix.lower()

            if ext == ".pdf":
                import pypdf
                reader = pypdf.PdfReader(str(path))
                text_pages = []
                for idx, page in enumerate(reader.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        text_pages.append(f"--- [第 {idx} 页] ---\n{text.strip()}")
                return "\n\n".join(text_pages) or "PDF 文件没有抓取到可提取文本内容。"

            elif ext == ".docx":
                import docx
                doc = docx.Document(str(path))
                lines = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        lines.append(p.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        lines.append(" | ".join(row_cells))
                return "\n".join(lines) or "Word 文档没有抓取到文本内容。"

            elif ext in (".pptx", ".ppt"):
                import pptx
                prs = pptx.Presentation(str(path))
                slides_text = []
                for idx, slide in enumerate(prs.slides, 1):
                    slide_lines = [f"=== [幻灯片 {idx}] ==="]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_lines.append(shape.text.strip())
                    if hasattr(slide, "notes_slide") and slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            slide_lines.append(f"[备注]: {notes}")
                    slides_text.append("\n".join(slide_lines))
                return "\n\n".join(slides_text) or "PPT 演示文稿没有抓取到文本内容。"

            elif ext in (".xlsx", ".xls", ".csv"):
                import pandas as pd
                def _df_to_text(df):
                    try:
                        return df.to_markdown(index=False)
                    except Exception:
                        return df.to_string(index=False)

                if ext == ".csv":
                    df = pd.read_csv(str(path))
                    return "=== [CSV 表格内容] ===\n" + _df_to_text(df)
                else:
                    excel_file = pd.ExcelFile(str(path))
                    sheet_contents = []
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(excel_file, sheet_name=sheet_name)
                        sheet_contents.append(f"=== [工作表: {sheet_name}] ===\n" + _df_to_text(df))
                    return "\n\n".join(sheet_contents) or "Excel 文件内容为空。"


            else:
                return f"暂不支持的文件扩展名: {ext}。目前支持 .pdf, .docx, .pptx, .xlsx, .csv。"

        except ValueError:
            return "拒绝访问：只能读取 workspace 文件夹内的内容。"
        except Exception as e:
            return f"解析办公文档失败: {e}"

    def create_word_document(self, relative_path: str, title: str, sections: list[dict[str, Any]]) -> str:
        """生成 Word (.docx) 文档。"""
        try:
            import docx
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            path = self._safe_path(relative_path)
            if not path.name.endswith(".docx"):
                path = path.with_suffix(".docx")
            path.parent.mkdir(parents=True, exist_ok=True)

            doc = docx.Document()
            if title:
                heading = doc.add_heading(title, level=0)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

            for sec in sections:
                sec_type = sec.get("type", "paragraph")
                if sec_type == "heading":
                    level = sec.get("level", 1)
                    doc.add_heading(sec.get("text", ""), level=level)
                elif sec_type == "paragraph":
                    doc.add_paragraph(sec.get("text", ""))
                elif sec_type == "bullet":
                    for item in sec.get("items", []):
                        doc.add_paragraph(item, style='List Bullet')
                elif sec_type == "table":
                    headers = sec.get("headers", [])
                    rows = sec.get("rows", [])
                    if headers or rows:
                        table = doc.add_table(rows=1 if headers else 0, cols=len(headers) if headers else (len(rows[0]) if rows else 1))
                        table.style = 'Table Grid'
                        if headers:
                            hdr_cells = table.rows[0].cells
                            for i, h in enumerate(headers):
                                hdr_cells[i].text = str(h)
                        for r in rows:
                            row_cells = table.add_row().cells
                            for i, val in enumerate(r):
                                if i < len(row_cells):
                                    row_cells[i].text = str(val)

            doc.save(str(path))
            return f"成功生成 Word 文档: {path.relative_to(self.workspace)}"
        except Exception as e:
            return f"生成 Word 文档失败: {e}"

    def create_ppt_presentation(self, relative_path: str, title: str, slides: list[dict[str, Any]]) -> str:
        """生成 PPT (.pptx) 演示文稿。"""
        try:
            import pptx

            path = self._safe_path(relative_path)
            if not path.name.endswith(".pptx"):
                path = path.with_suffix(".pptx")
            path.parent.mkdir(parents=True, exist_ok=True)

            prs = pptx.Presentation()

            # 封面页
            if title:
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                slide.shapes.title.text = title

            # 内容页
            bullet_slide_layout = prs.slide_layouts[1]
            for s_data in slides:
                slide = prs.slides.add_slide(bullet_slide_layout)
                s_title = s_data.get("title", "")
                bullets = s_data.get("bullets", [])
                notes = s_data.get("notes", "")

                if s_title and slide.shapes.title:
                    slide.shapes.title.text = s_title

                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()

                for i, b_text in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = b_text
                    p.level = 0

                if notes and hasattr(slide, "notes_slide") and slide.notes_slide:
                    slide.notes_slide.notes_text_frame.text = notes

            prs.save(str(path))
            return f"成功生成 PPT 演示文稿: {path.relative_to(self.workspace)}"
        except Exception as e:
            return f"生成 PPT 失败: {e}"

    def create_pdf_document(self, relative_path: str, title: str, content_markdown: str) -> str:
        """将 Markdown / HTML 内容高质量导出为 PDF 文件。"""
        try:
            path = self._safe_path(relative_path)
            if not path.name.endswith(".pdf"):
                path = path.with_suffix(".pdf")
            path.parent.mkdir(parents=True, exist_ok=True)

            import html
            formatted_content = content_markdown.replace("\n", "<br/>")
            html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{html.escape(title)}</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif; margin: 40px; line-height: 1.6; color: #333; }}
        h1 {{ border-bottom: 2px solid #0066cc; padding-bottom: 10px; color: #111; }}
        code {{ background: #f4f4f4; padding: 2px 5px; border-radius: 3px; font-family: monospace; }}
    </style>
</head>
<body>
    <h1>{html.escape(title)}</h1>
    <div>{formatted_content}</div>
</body>
</html>"""

            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page()
                page.set_content(html_content, wait_until="networkidle")
                page.pdf(path=str(path), format="A4", print_background=True, margin={"top": "20mm", "bottom": "20mm", "left": "20mm", "right": "20mm"})
                browser.close()

            return f"成功生成 PDF 文档: {path.relative_to(self.workspace)}"
        except Exception as e:
            return f"导出 PDF 文档失败: {e}"

    def create_excel_spreadsheet(self, relative_path: str, sheets_data: dict[str, list[list[Any]]]) -> str:
        """生成 Excel (.xlsx) 表格文件。"""
        try:
            import pandas as pd

            path = self._safe_path(relative_path)
            if not path.name.endswith(".xlsx"):
                path = path.with_suffix(".xlsx")
            path.parent.mkdir(parents=True, exist_ok=True)

            with pd.ExcelWriter(str(path), engine='openpyxl') as writer:
                for sheet_name, data in sheets_data.items():
                    if not data:
                        df = pd.DataFrame()
                    elif len(data) == 1:
                        df = pd.DataFrame(data)
                    else:
                        headers = data[0]
                        rows = data[1:]
                        df = pd.DataFrame(rows, columns=headers)
                    df.to_excel(writer, sheet_name=sheet_name[:31], index=False)

            return f"成功生成 Excel 表格文件: {path.relative_to(self.workspace)}"
        except Exception as e:
            return f"生成 Excel 文件失败: {e}"

    # --- 时间与环境信息 ---

    def get_current_datetime(self, timezone: str = "Asia/Shanghai") -> str:
        """获取当前系统日期、时间、星期与时区信息。"""
        import datetime
        try:
            import zoneinfo
            tz = zoneinfo.ZoneInfo(timezone)
        except Exception:
            tz = None
        now = datetime.datetime.now(tz) if tz else datetime.datetime.now()
        weekdays = ["周一", "周二", "周三", "周四", "周五", "周六", "周日"]
        return (
            f"当前时间：{now.strftime('%Y年%m月%d日 %H:%M:%S')} "
            f"{weekdays[now.weekday()]}，时区：{timezone}"
        )

    # --- 框架级通用基础工具能力 ---

    def search_files(self, pattern: str, relative_path: str = ".", file_glob: str = "*") -> str:
        """在 workspace 指定目录下递归搜索包含特定文本模式的文件内容 (grep)。"""
        try:
            target_dir = self._safe_path(relative_path or ".")
            if not target_dir.exists():
                return f"路径不存在: {relative_path}"
            
            matches = []
            max_matches = 100
            count = 0
            
            search_items = target_dir.rglob(file_glob) if target_dir.is_dir() else [target_dir]
            for file_path in search_items:
                if not file_path.is_file():
                    continue
                if file_path.name.startswith(".") or ".git" in file_path.parts or "__pycache__" in file_path.parts:
                    continue
                try:
                    lines = file_path.read_text(encoding="utf-8", errors="ignore").splitlines()
                    for idx, line in enumerate(lines, start=1):
                        if pattern.lower() in line.lower():
                            rel = file_path.relative_to(self.workspace)
                            matches.append(f"{rel}:{idx}: {line.strip()[:150]}")
                            count += 1
                            if count >= max_matches:
                                break
                except Exception:
                    continue
                if count >= max_matches:
                    break
            
            if not matches:
                return f"未在 {relative_path} 下找到包含 '{pattern}' 的文本内容。"
            res = "\n".join(matches)
            if count >= max_matches:
                res += f"\n\n...(搜索结果过多，仅显示前 {max_matches} 条)"
            return res
        except Exception as e:
            return f"搜索文件内容失败: {e}"

    def http_request(self, url: str, method: str = "GET", headers: dict = None, body: str = None, timeout: float = 15.0) -> str:
        """发起通用 HTTP/REST API 请求 (GET/POST/PUT/DELETE)。"""
        if not url.startswith(("http://", "https://")):
            return "URL 格式无效，必须以 http:// 或 https:// 开头。"
        try:
            req_headers = headers or {}
            if "User-Agent" not in req_headers:
                req_headers["User-Agent"] = "DamingAgent/1.0"
            
            content = body.encode("utf-8") if isinstance(body, str) else None
            response = httpx.request(
                method=method.upper(),
                url=url,
                headers=req_headers,
                content=content,
                timeout=float(timeout),
                follow_redirects=True,
            )
            res_text = response.text
            if len(res_text) > 4000:
                res_text = res_text[:4000] + "\n...(响应内容过长已截断)"
            return f"状态码: {response.status_code}\n响应内容:\n{res_text}"
        except Exception as e:
            return f"HTTP 请求失败: {e}"

    def clarify(self, question: str) -> str:
        """向用户提出结构化澄清问题。"""
        return f"❓ [Agent 澄清提问]: {question}"

    def manage_plan(self, action: str, title: str = "", steps: list[dict[str, Any]] = None, step_id: int = None, status: str = None) -> str:
        """动态任务计划管理器 (create / get / update_step / list)。"""
        plan_file = self.workspace / ".runtime_plan.json"
        try:
            plan_data = {}
            if plan_file.exists():
                try:
                    plan_data = json.loads(plan_file.read_text(encoding="utf-8"))
                except Exception:
                    plan_data = {}
            
            act = action.lower().strip()
            if act == "create":
                plan_data = {
                    "title": title or "未命名计划",
                    "updated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
                    "steps": []
                }
                if steps:
                    for i, s in enumerate(steps, 1):
                        plan_data["steps"].append({
                            "id": s.get("id", i),
                            "description": s.get("description", ""),
                            "status": s.get("status", "pending"),
                            "depends_on": s.get("depends_on", [])
                        })
                plan_file.write_text(json.dumps(plan_data, ensure_ascii=False, indent=2), encoding="utf-8")
                return f"已成功创建执行计划: 【{plan_data['title']}】，共 {len(plan_data['steps'])} 个步骤。"

            elif act == "update_step":
                if not plan_data:
                    return "当前没有活跃的计划，请先创建计划 (action='create')。"
                if step_id is None:
                    return "更新步骤需指定 step_id。"
                found = False
                for s in plan_data.get("steps", []):
                    if s.get("id") == step_id:
                        if status:
                            s["status"] = status
                        found = True
                        break
                if not found:
                    return f"计划中未找到步骤 ID: {step_id}"
                plan_data["updated_at"] = time.strftime("%Y-%m-%d %H:%M:%S")
                plan_file.write_text(json.dumps(plan_data, ensure_ascii=False, indent=2), encoding="utf-8")
                return f"已将计划步骤 [{step_id}] 状态更新为: {status}"

            elif act in ("get", "list", "show"):
                if not plan_data:
                    return "当前没有活跃的计划。"
                lines = [f"📋 [执行计划]: {plan_data.get('title')} (更新时间: {plan_data.get('updated_at')})"]
                status_symbols = {"pending": "⏳ 待执行", "in_progress": "🔄 执行中", "completed": "✅ 已完成", "failed": "❌ 失败"}
                for s in plan_data.get("steps", []):
                    sym = status_symbols.get(s.get("status"), s.get("status"))
                    deps = f" (依赖: {s['depends_on']})" if s.get("depends_on") else ""
                    lines.append(f"  [{s.get('id')}] {sym} - {s.get('description')}{deps}")
                return "\n".join(lines)

            else:
                return f"未知 action: {action}。支持: create, update_step, get。"
        except Exception as e:
            return f"操作计划管理失败: {e}"

    def analyze_image(self, image_path: str, prompt: str = "请详细描述这张图片的内容") -> str:
        """分析本地图像或截图内容。"""
        try:
            path = self._safe_path(image_path)
            if not path.is_file():
                return f"图片文件不存在: {image_path}"
            return f"📷 [已载入图片]: {path.name} (大小: {path.stat().st_size} bytes)。提示词: {prompt}。\n(注: 图片已准备就绪，包含图片输入上下文。)"
        except Exception as e:
            return f"读取分析图片失败: {e}"

    def read_clipboard(self) -> str:
        """读取系统剪贴板文本内容。"""
        try:
            res = subprocess.run(["pbpaste"], capture_output=True, text=True, timeout=5)
            text = res.stdout.strip()
            return text or "剪贴板内容为空。"
        except Exception as e:
            return f"读取剪贴板失败: {e}"

    def write_clipboard(self, text: str) -> str:
        """写入文本到系统剪贴板。"""
        try:
            p = subprocess.Popen(["pbcopy"], stdin=subprocess.PIPE, text=True)
            p.communicate(input=text)
            return "成功将内容写入系统剪贴板。"
        except Exception as e:
            return f"写入剪贴板失败: {e}"

    def notify(self, title: str, message: str) -> str:
        """发送 Mac 系统桌面通知。"""
        try:
            script = f'display notification "{message}" with title "{title}"'
            subprocess.run(["osascript", "-e", script], check=True, timeout=5)
            return f"已成功发送系统通知: 【{title}】"
        except Exception as e:
            return f"发送系统通知失败: {e}"

    def calculate(self, expression: str) -> str:
        """安全计算数学表达式 (例如: '2**10 + 50 * 3')。"""
        try:
            allowed_names = {"abs": abs, "round": round, "min": min, "max": max, "sum": sum, "pow": pow}
            # 使用 compile AST 校验安全表达式，杜绝 __import__ 等注入
            code = compile(expression, "<string>", "eval")
            for name in code.co_names:
                if name not in allowed_names:
                    return f"拒绝计算：表达式包含不安全的变量/函数引用 '{name}'。"
            result = eval(code, {"__builtins__": {}}, allowed_names)
            return f"计算结果: {result}"
        except Exception as e:
            return f"计算数学表达式失败: {e}"

    # --- 工具助手 ---

    def _safe_path(self, relative_path: str) -> Path:
        path = (self.workspace / relative_path).resolve()
        if self.workspace != path and self.workspace not in path.parents:
            raise ValueError("path outside workspace")
        return path
